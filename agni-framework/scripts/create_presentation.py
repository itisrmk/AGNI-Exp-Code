#!/usr/bin/env python3
"""
Generate AGNI Framework Presentation
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0x1A, 0x56, 0x8F)  # Dark blue
ACCENT_COLOR = RGBColor(0x2E, 0xCC, 0x71)  # Green
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)  # Dark gray
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12.333)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(left, Inches(4.2), width, Inches(1.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.4), Inches(12), Inches(5.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(point, tuple):
            # Main point with sub-points
            p.text = "• " + point[0]
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_after = Pt(8)

            for sub in point[1]:
                p = tf.add_paragraph()
                p.text = "    ‣ " + sub
                p.font.size = Pt(18)
                p.font.color.rgb = LIGHT_GRAY
                p.level = 1
                p.space_after = Pt(4)
        else:
            p.text = "• " + point
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_after = Pt(8)

    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(12.333)
    height = Inches(0.5 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set column widths
    col_width = Inches(12.333 / num_cols)
    for col in table.columns:
        col.width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


# ============== SLIDES ==============

# Slide 1: Title
add_title_slide(
    prs,
    "AGNI: Adaptive Glucose Neural Intelligence",
    "A Framework for Continual Learning in Continuous Glucose Monitoring\n\nDoctoral Dissertation Research",
)

# Slide 2: The Problem
add_content_slide(
    prs,
    "The Problem: Why This Research Matters",
    [
        "Diabetes affects 537 million people worldwide",
        (
            "Continuous Glucose Monitors (CGMs) help manage blood sugar",
            [
                "Measure glucose every 5 minutes",
                "Can predict future glucose levels",
                "Early warnings prevent dangerous highs and lows",
            ],
        ),
        (
            "The Challenge: Individual bodies change over time",
            [
                "Lifestyle changes, stress, illness affect patterns",
                "Static prediction models become outdated",
                "Need models that adapt to each person",
            ],
        ),
        "Our Goal: Create AI that learns and adapts continuously",
    ],
)

# Slide 3: What is Continual Learning?
add_content_slide(
    prs,
    "What is Continual Learning?",
    [
        (
            "Traditional Machine Learning:",
            [
                "Train once, deploy forever",
                "Like memorizing a textbook - doesn't update",
                "Performance degrades as patterns change",
            ],
        ),
        (
            "Continual Learning:",
            [
                "Learns from new data while retaining old knowledge",
                "Like a doctor who learns from each new patient",
                "Adapts to individual changes over time",
            ],
        ),
        "Key Challenge: Avoid 'catastrophic forgetting' - learning new things without forgetting important old patterns",
    ],
)

# Slide 4: Research Approach - Three Strategies
add_content_slide(
    prs,
    "Our Research: Comparing Three Adaptation Strategies",
    [
        (
            "Strategy 1: Static (Baseline)",
            ["Train once, never update", "Simple but degrades over time"],
        ),
        (
            "Strategy 2: Periodic Retraining",
            [
                "Retrain from scratch every week",
                "Fresh model but computationally expensive",
            ],
        ),
        (
            "Strategy 3: Continual Learning (Our Innovation)",
            [
                "Daily micro-updates with smart techniques",
                "EWC: Protects important learned patterns",
                "Experience Replay: Remembers critical events",
            ],
        ),
        "We compare all three using real patient data",
    ],
)

# Slide 5: The Dataset
add_content_slide(
    prs,
    "OhioT1DM Dataset: Real Patient Data",
    [
        "12 patients with Type 1 Diabetes",
        "8 weeks of continuous monitoring per patient",
        (
            "Data collected every 5 minutes:",
            [
                "Blood glucose readings (CGM)",
                "Meal carbohydrate intake",
                "Insulin doses",
                "Exercise, sleep, stress events",
            ],
        ),
        "~16,000 glucose readings per patient",
        "Benchmark dataset used by researchers worldwide",
    ],
)

# Slide 6: Project Stages
add_content_slide(
    prs,
    "Project Stages: Building Block by Block",
    [
        (
            "STAGE 1: Foundation & Static Baseline (COMPLETED)",
            [
                "Build data pipeline and model architectures",
                "Establish baseline performance to beat",
            ],
        ),
        (
            "STAGE 2: Periodic Retraining (Next)",
            ["Implement weekly retraining strategy", "Compare against static baseline"],
        ),
        (
            "STAGE 3: Continual Learning (Final)",
            [
                "Implement EWC + Experience Replay",
                "Run ablation studies to prove each component works",
                "Statistical analysis and publication",
            ],
        ),
    ],
)

# Slide 7: Stage 1 - What We Built
add_content_slide(
    prs,
    "Stage 1: What We Built",
    [
        (
            "Data Pipeline:",
            [
                "XML parser for OhioT1DM dataset",
                "Preprocessing: gap handling, normalization",
                "PyTorch datasets with sliding windows",
            ],
        ),
        (
            "Three Neural Network Models:",
            [
                "LSTM: Sequential memory network (50K parameters)",
                "TCN: Temporal convolutions (22K parameters)",
                "Transformer: Attention-based (34K parameters)",
            ],
        ),
        (
            "Evaluation Framework:",
            [
                "Standard metrics: MAE, RMSE, Correlation",
                "Clinical metrics: Clarke Error Grid",
                "Hypoglycemia detection analysis",
            ],
        ),
    ],
)

# Slide 8: The Three Models Explained
add_two_column_slide(
    prs,
    "Understanding Our Three Models",
    [
        "LSTM (Long Short-Term Memory)",
        "• Remembers patterns over time",
        "• Like reading a book sequentially",
        "• Good at capturing trends",
        "",
        "TCN (Temporal Convolutional Network)",
        "• Uses filters to detect patterns",
        "• Like scanning for specific shapes",
        "• Fast and efficient",
    ],
    [
        "Transformer",
        "• Pays attention to important moments",
        "• Can look at all history at once",
        "• State-of-the-art architecture",
        "",
        "All predict glucose 30 minutes ahead",
        "using the last 2 hours of data",
        "(24 readings at 5-min intervals)",
    ],
)

# Slide 9: Results Table
add_table_slide(
    prs,
    "Stage 1 Results: Static Baseline (30-min Prediction)",
    ["Model", "MAE (mg/dL)", "RMSE (mg/dL)", "Correlation", "Clarke A+B"],
    [
        ["Transformer", "18.26", "24.99", "0.89", "97.52%"],
        ["LSTM", "18.80", "25.81", "0.88", "97.68%"],
        ["TCN", "20.18", "26.83", "0.87", "96.61%"],
    ],
)

# Slide 10: Understanding the Metrics
add_content_slide(
    prs,
    "What Do These Numbers Mean?",
    [
        (
            "MAE (Mean Absolute Error): 18-20 mg/dL",
            [
                "Average prediction error",
                "Lower is better",
                "State-of-the-art is ~13-15 mg/dL",
            ],
        ),
        (
            "RMSE: 25-27 mg/dL",
            [
                "Penalizes large errors more heavily",
                "Important for safety-critical predictions",
            ],
        ),
        (
            "Clarke A+B: 96-98%",
            [
                "Clinical safety metric",
                "% of predictions in safe zones",
                ">95% considered clinically acceptable",
            ],
        ),
        "All models achieve clinically acceptable performance!",
    ],
)

# Slide 11: Per-Patient Results
add_table_slide(
    prs,
    "Per-Patient Results (Transformer Model)",
    ["Patient", "MAE", "RMSE", "Correlation", "Time in Range"],
    [
        ["559", "21.34", "29.55", "0.920", "56.2%"],
        ["563", "14.65", "19.39", "0.893", "74.2%"],
        ["570", "13.66", "18.36", "0.955", "43.2%"],
        ["575", "15.98", "20.82", "0.914", "69.2%"],
        ["588", "20.34", "27.12", "0.888", "63.7%"],
        ["591", "23.58", "34.68", "0.758", "64.2%"],
    ],
)

# Slide 12: Key Findings
add_content_slide(
    prs,
    "Key Findings from Stage 1",
    [
        "Transformer achieves best overall accuracy",
        "LSTM achieves best clinical safety (Clarke A+B)",
        "TCN most parameter-efficient (smallest model)",
        (
            "Patient variability is significant:",
            [
                "Best patient (570): MAE = 13.66 mg/dL",
                "Worst patient (591): MAE = 23.58 mg/dL",
                "This is why personalized adaptation matters!",
            ],
        ),
        "All models provide clinically acceptable predictions",
        "Static baseline established - ready for adaptation strategies",
    ],
)

# Slide 13: What's Next - Stage 2
add_content_slide(
    prs,
    "Next: Stage 2 - Periodic Retraining",
    [
        (
            "What we'll implement:",
            [
                "Weekly model retraining from scratch",
                "Use most recent 7 days of data",
                "Compare against static baseline",
            ],
        ),
        (
            "Expected outcome:",
            [
                "Sawtooth pattern: good after retraining, degrades over week",
                "Better than static on average",
                "But computationally expensive",
            ],
        ),
        "This sets up the comparison for Stage 3's continual learning",
    ],
)

# Slide 14: Stage 3 Preview
add_content_slide(
    prs,
    "Final Stage: Continual Learning",
    [
        (
            "Elastic Weight Consolidation (EWC):",
            [
                "Identifies which model parameters are 'important'",
                "Protects them during updates",
                "Prevents catastrophic forgetting",
            ],
        ),
        (
            "Experience Replay Buffer:",
            [
                "Stores critical past examples",
                "Prioritizes rare events (low/high glucose)",
                "Mixes old and new data during training",
            ],
        ),
        (
            "Ablation Studies:",
            [
                "Test each component individually",
                "Prove both are necessary",
                "Optimize buffer size and update frequency",
            ],
        ),
    ],
)

# Slide 15: Expected Impact
add_content_slide(
    prs,
    "Expected Results & Impact",
    [
        (
            "Expected Performance Improvement:",
            [
                "Static: ~18-20 mg/dL MAE (degrades over time)",
                "Periodic: ~15-17 mg/dL MAE (sawtooth pattern)",
                "Continual: ~13-15 mg/dL MAE (stable performance)",
            ],
        ),
        (
            "Clinical Impact:",
            [
                "Better hypoglycemia detection saves lives",
                "Personalized predictions for each patient",
                "Models that improve over time, not degrade",
            ],
        ),
        (
            "Research Contribution:",
            [
                "First comprehensive comparison of adaptation strategies",
                "Open-source framework for CGM research",
                "Reproducible benchmark results",
            ],
        ),
    ],
)

# Slide 16: Technical Architecture
add_content_slide(
    prs,
    "Technical Implementation",
    [
        (
            "Platform:",
            [
                "Python 3.9 with PyTorch 2.8",
                "Apple Silicon (MPS) GPU acceleration",
                "Modular, extensible codebase",
            ],
        ),
        (
            "Project Structure:",
            [
                "src/data/ - Data loading and preprocessing",
                "src/models/ - LSTM, TCN, Transformer",
                "src/adaptation/ - Adaptation strategies",
                "src/evaluation/ - Metrics and visualization",
            ],
        ),
        "~3,000 lines of well-documented code",
        "Configuration-driven experiments for reproducibility",
    ],
)

# Slide 17: Summary
add_content_slide(
    prs,
    "Summary",
    [
        "Problem: Static glucose prediction models degrade over time",
        "Solution: Continual learning that adapts to each patient",
        (
            "Stage 1 Complete:",
            [
                "Built complete framework with 3 model architectures",
                "Established baseline: 18.26 mg/dL MAE (Transformer)",
                "All models achieve >96% clinical acceptability",
            ],
        ),
        (
            "Coming Next:",
            [
                "Stage 2: Periodic retraining comparison",
                "Stage 3: Continual learning with EWC + Replay",
            ],
        ),
        "Goal: Personalized AI that learns and adapts like a doctor",
    ],
)

# Slide 18: Thank You
add_title_slide(
    prs,
    "Thank You",
    "Questions?\n\nAGNI: Adaptive Glucose Neural Intelligence Framework",
)

# Save presentation
output_path = (
    "/Users/rahulkashyap/Desktop/DIT/Presentation/AGNI_Framework_Presentation.pptx"
)
output_path = "/Users/rahulkashyap/Desktop/DIT/Presentation/AGNI_Framework_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
