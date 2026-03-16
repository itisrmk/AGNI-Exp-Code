#!/usr/bin/env python3
"""
Create Stage 2 PowerPoint Presentation for AGNI Framework
Periodic Retraining Results and Analysis
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4), Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()

        para.text = f"• {point}"
        para.font.size = Pt(20)
        para.space_after = Pt(12)
        para.font.color.rgb = RGBColor(50, 50, 50)

    return slide


def add_table_slide(prs, title, headers, data, highlight_row=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Table
    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * rows)
    ).table

    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(9 / cols)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER

            if highlight_row is not None and row_idx == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 230, 200)

    return slide


def add_comparison_slide(prs, title, comparison_data):
    """Add a comparison slide with static vs periodic."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Create comparison boxes
    y_pos = 1.5
    for model, static_mae, periodic_mae, delta in comparison_data:
        # Model name
        model_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5)
        )
        model_frame = model_box.text_frame
        model_para = model_frame.paragraphs[0]
        model_para.text = model
        model_para.font.size = Pt(20)
        model_para.font.bold = True

        # Static box
        static_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5),
            Inches(y_pos),
            Inches(2.5),
            Inches(0.6),
        )
        static_shape.fill.solid()
        static_shape.fill.fore_color.rgb = RGBColor(100, 149, 237)
        static_shape.text_frame.paragraphs[0].text = f"Static: {static_mae}"
        static_shape.text_frame.paragraphs[0].font.size = Pt(14)
        static_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        static_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Periodic box
        periodic_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2),
            Inches(y_pos),
            Inches(2.5),
            Inches(0.6),
        )
        periodic_shape.fill.solid()
        if delta > 0:
            periodic_shape.fill.fore_color.rgb = RGBColor(
                50, 205, 50
            )  # Green - improvement
        elif delta < -5:
            periodic_shape.fill.fore_color.rgb = RGBColor(
                220, 20, 60
            )  # Red - degradation
        else:
            periodic_shape.fill.fore_color.rgb = RGBColor(
                255, 165, 0
            )  # Orange - similar
        periodic_shape.text_frame.paragraphs[0].text = f"Periodic: {periodic_mae}"
        periodic_shape.text_frame.paragraphs[0].font.size = Pt(14)
        periodic_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        periodic_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Delta
        delta_box = slide.shapes.add_textbox(
            Inches(7.9), Inches(y_pos), Inches(1.5), Inches(0.6)
        )
        delta_frame = delta_box.text_frame
        delta_para = delta_frame.paragraphs[0]
        delta_para.text = f"Δ {delta:+.2f}"
        delta_para.font.size = Pt(16)
        delta_para.font.bold = True
        if delta > 0:
            delta_para.font.color.rgb = RGBColor(0, 128, 0)
        elif delta < -5:
            delta_para.font.color.rgb = RGBColor(220, 20, 60)
        else:
            delta_para.font.color.rgb = RGBColor(255, 165, 0)

        y_pos += 1.2

    return slide


def create_presentation():
    """Create the complete Stage 2 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs, "AGNI Framework - Stage 2", "Periodic Retraining Strategy & Results"
    )

    # Slide 2: Recap - What We Did in Stage 1
    add_content_slide(
        prs,
        "Stage 1 Recap: Static Baseline",
        [
            "Built complete glucose prediction framework with 3 neural network models",
            "LSTM: Recurrent architecture with memory (50K parameters)",
            "TCN: Convolutional approach with dilated filters (22K parameters)",
            "Transformer: Attention-based architecture (34K parameters)",
            "Trained once on initial data, never updated",
            "Best result: Transformer with MAE 18.26 mg/dL",
            "All models achieved >96% clinical acceptability (Clarke A+B)",
        ],
    )

    # Slide 3: The Problem with Static Models
    add_content_slide(
        prs,
        "The Problem: Why Static Models Degrade",
        [
            "Patient physiology changes over time (diet, exercise, stress, illness)",
            "Glucose patterns that worked last month may not work today",
            "Static models are frozen in time - they can't adapt",
            "Performance gradually degrades as patterns drift",
            "Question: Can we improve by periodically updating the model?",
        ],
    )

    # Slide 4: Stage 2 Goal
    add_content_slide(
        prs,
        "Stage 2 Goal: Periodic Retraining",
        [
            "Strategy: Retrain the model from scratch at fixed intervals",
            "Implementation: Every 7 days, create a fresh model",
            "Training data: Most recent 14 days (sliding window)",
            "Hypothesis: Fresh models trained on recent data should perform better",
            "This is a common industry approach - regularly refresh models",
        ],
    )

    # Slide 5: How Periodic Retraining Works
    add_content_slide(
        prs,
        "How Periodic Retraining Works",
        [
            "Day 1-7: Initial training on first week of data",
            "Day 8-14: Model deployed, evaluated daily",
            "Day 14: RETRAIN - throw away old model, train new one",
            "Day 15-21: New model deployed, evaluated daily",
            "Day 21: RETRAIN again... and so on",
            "Expected pattern: 'Sawtooth' - best after retrain, degrades over time",
        ],
    )

    # Slide 6: Implementation Details
    add_content_slide(
        prs,
        "Implementation Details",
        [
            "Created PeriodicAdapter class in src/adaptation/periodic.py",
            "Retraining interval: 7 days",
            "Training window: 14 days of recent data",
            "Same hyperparameters as Stage 1 (50 epochs, early stopping)",
            "Evaluation: Daily metrics + Clarke Error Grid",
            "Ran experiments on all 6 patients × 3 models = 18 experiments",
        ],
    )

    # Slide 7: LSTM Results Table
    add_table_slide(
        prs,
        "LSTM Periodic Results",
        ["Patient", "MAE (mg/dL)", "RMSE (mg/dL)", "Clarke A+B", "Retrains"],
        [
            ["559", "38.22", "46.59", "97.6%", "4"],
            ["563", "25.74", "31.99", "98.1%", "5"],
            ["570", "32.78", "39.90", "97.6%", "4"],
            ["575", "33.95", "41.71", "96.6%", "5"],
            ["588", "29.83", "36.65", "97.2%", "5"],
            ["591", "33.68", "40.49", "96.4%", "5"],
            ["MEAN", "32.37", "39.56", "97.2%", "-"],
        ],
        highlight_row=6,
    )

    # Slide 8: TCN Results Table
    add_table_slide(
        prs,
        "TCN Periodic Results",
        ["Patient", "MAE (mg/dL)", "RMSE (mg/dL)", "Clarke A+B", "Retrains"],
        [
            ["559", "22.05", "30.55", "95.9%", "4"],
            ["563", "15.17", "20.62", "98.5%", "5"],
            ["570", "15.43", "20.78", "98.5%", "4"],
            ["575", "19.73", "27.10", "96.1%", "5"],
            ["588", "18.81", "24.83", "96.6%", "5"],
            ["591", "23.96", "31.65", "95.9%", "5"],
            ["MEAN", "19.19", "25.92", "96.9%", "-"],
        ],
        highlight_row=6,
    )

    # Slide 9: Transformer Results Table
    add_table_slide(
        prs,
        "Transformer Periodic Results",
        ["Patient", "MAE (mg/dL)", "RMSE (mg/dL)", "Clarke A+B", "Retrains"],
        [
            ["559", "20.81", "29.01", "96.3%", "4"],
            ["563", "14.38", "19.93", "98.7%", "5"],
            ["570", "14.89", "20.48", "98.3%", "4"],
            ["575", "19.78", "27.47", "95.5%", "5"],
            ["588", "18.12", "24.20", "97.6%", "5"],
            ["591", "21.39", "28.82", "96.3%", "5"],
            ["MEAN", "18.23", "24.99", "97.1%", "-"],
        ],
        highlight_row=6,
    )

    # Slide 10: Static vs Periodic Comparison
    add_comparison_slide(
        prs,
        "Static vs Periodic: MAE Comparison",
        [
            ("LSTM", "18.80", "32.37", -13.57),
            ("TCN", "20.18", "19.19", 0.99),
            ("Transformer", "18.26", "18.23", 0.03),
        ],
    )

    # Slide 11: Complete Comparison Table
    add_table_slide(
        prs,
        "Complete Strategy Comparison",
        ["Model", "Strategy", "MAE", "RMSE", "Clarke A+B"],
        [
            ["LSTM", "Static", "18.80", "25.81", "97.7%"],
            ["LSTM", "Periodic", "32.37", "39.56", "97.2%"],
            ["TCN", "Static", "20.18", "26.83", "96.6%"],
            ["TCN", "Periodic", "19.19", "25.92", "96.9%"],
            ["Transformer", "Static", "18.26", "24.99", "97.5%"],
            ["Transformer", "Periodic", "18.23", "24.99", "97.1%"],
        ],
    )

    # Slide 12: Key Finding 1 - LSTM Degradation
    add_content_slide(
        prs,
        "Key Finding 1: LSTM Performance Degradation",
        [
            "LSTM MAE increased from 18.80 to 32.37 mg/dL (-72% worse!)",
            "Why? LSTM learns long-term temporal patterns",
            "These patterns require extensive historical data to learn",
            "When we retrain with only 14 days, we lose this knowledge",
            "LSTM essentially 'forgets' the important patterns it learned",
            "Conclusion: Periodic retraining hurts LSTM significantly",
        ],
    )

    # Slide 13: Key Finding 2 - TCN/Transformer Robustness
    add_content_slide(
        prs,
        "Key Finding 2: TCN & Transformer Robustness",
        [
            "TCN improved slightly: 20.18 → 19.19 mg/dL (+5% better)",
            "Transformer nearly identical: 18.26 → 18.23 mg/dL",
            "Why are these models more robust?",
            "TCN: Convolutional filters learn local patterns efficiently",
            "Transformer: Attention mechanism focuses on relevant features",
            "Both architectures need less data to learn effective patterns",
            "Conclusion: Architecture matters for adaptation strategies!",
        ],
    )

    # Slide 14: Key Finding 3 - Clinical Safety Maintained
    add_content_slide(
        prs,
        "Key Finding 3: Clinical Safety Maintained",
        [
            "All models maintained >96% Clarke A+B across both strategies",
            "This is the clinical threshold for acceptable glucose prediction",
            "Even LSTM with higher MAE still meets clinical requirements",
            "Periodic retraining doesn't compromise safety",
            "Important: Accuracy metrics alone don't tell the full story",
            "Clinical metrics ensure predictions are safe for patient use",
        ],
    )

    # Slide 15: The Sawtooth Pattern
    add_content_slide(
        prs,
        "Observed: The Sawtooth Pattern",
        [
            "As expected, performance follows a sawtooth curve:",
            "Day 1 after retrain: Best performance (fresh model)",
            "Days 2-7: Gradual degradation as patterns drift",
            "Day 7: Performance at its worst",
            "Retrain: Performance jumps back up",
            "This pattern repeats every retraining cycle",
            "Indicates the model is indeed adapting to recent data",
        ],
    )

    # Slide 16: What We Learned
    add_content_slide(
        prs,
        "What We Learned from Stage 2",
        [
            "Periodic retraining is NOT universally beneficial",
            "LSTM: Requires all historical data - don't retrain from scratch",
            "TCN/Transformer: More robust to limited training windows",
            "Model architecture determines optimal adaptation strategy",
            "Clinical safety is maintained regardless of strategy",
            "Need a smarter approach: Continual Learning (Stage 3)",
        ],
    )

    # Slide 17: Preview of Stage 3
    add_content_slide(
        prs,
        "Preview: Stage 3 - Continual Learning",
        [
            "Goal: Get the best of both worlds",
            "Elastic Weight Consolidation (EWC): Protect important knowledge",
            "Experience Replay: Remember critical past events",
            "Daily micro-updates instead of weekly full retraining",
            "Hypothesis: Should outperform both Static and Periodic",
            "Especially important for LSTM which suffered with periodic retraining",
        ],
    )

    # Slide 18: Thank You
    add_title_slide(
        prs, "Questions?", "Stage 2 Complete - Ready for Stage 3: Continual Learning"
    )

    # Save
    output_path = Path(
        "/Users/rahulkashyap/Desktop/DIT/Presentation/AGNI_Stage2_Presentation.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
